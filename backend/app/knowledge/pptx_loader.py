"""
PPTX ve DOCX dosyalarından metin çıkarma ve parçalara ayırma.

İşlevler:
    extract_slides      – PPTX slaytlarından metin çıkarır.
    extract_docx_qa     – DOCX'ten Soru-Cevap çiftlerini çıkarır.
    chunk_text          – Metni belirli boyutta örtüşen parçalara böler.
    load_and_chunk_pptx – PPTX: çıkar + parçala.
    load_and_chunk_docx – DOCX: SSS çiftlerini chunk olarak döndürür.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import TypedDict

from pptx import Presentation

logger = logging.getLogger("ogrenci_destek.knowledge.loader")


# ── Tip tanımları ─────────────────────────────────────────────────────
class SlideText(TypedDict):
    slide_number: int
    text: str


class Chunk(TypedDict):
    text: str
    slide_number: int | None
    chunk_index: int
    source: str  # "pptx" | "docx"


# ── Beyaz boşluk normalleştirme ──────────────────────────────────────
_MULTI_SPACE = re.compile(r"[ \t]+")
_MULTI_NEWLINE = re.compile(r"\n{3,}")


def _normalize(raw: str) -> str:
    """Fazla boşlukları temizle, Türkçe karakterleri koru."""
    text = _MULTI_SPACE.sub(" ", raw)
    text = _MULTI_NEWLINE.sub("\n\n", text)
    return text.strip()


# ══════════════════════════════════════════════════════════════════════
#  PPTX İŞLEMLERİ
# ══════════════════════════════════════════════════════════════════════

def extract_slides(pptx_path: str | Path) -> list[SlideText]:
    """
    PPTX dosyasındaki tüm slaytlardan metin çıkarır.

    Args:
        pptx_path: .pptx dosyasının yolu.

    Returns:
        Her slayt için {slide_number, text} listesi.
    """
    pptx_path = Path(pptx_path)
    if not pptx_path.exists():
        raise FileNotFoundError(f"PPTX dosyası bulunamadı: {pptx_path}")

    prs = Presentation(str(pptx_path))
    slides: list[SlideText] = []

    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    line = paragraph.text.strip()
                    if line:
                        parts.append(line)

            # Tablolardan da metin çıkar
            if shape.has_table:
                for row in shape.table.rows:
                    row_texts = [
                        cell.text.strip()
                        for cell in row.cells
                        if cell.text.strip()
                    ]
                    if row_texts:
                        parts.append(" | ".join(row_texts))

        full_text = _normalize("\n".join(parts))
        if full_text:
            slides.append(SlideText(slide_number=idx, text=full_text))

    logger.info("PPTX'ten %d slayt çıkarıldı: %s", len(slides), pptx_path.name)
    return slides


# ══════════════════════════════════════════════════════════════════════
#  DOCX İŞLEMLERİ
# ══════════════════════════════════════════════════════════════════════

def extract_docx_qa(docx_path: str | Path) -> list[dict]:
    """
    DOCX dosyasından Soru-Cevap çiftlerini çıkarır.

    SSS formatını algılar – iki olası yapı desteklenir:
      A) Tek paragrafta: "Soru:…\\nCevap:…" (non-breaking space olabilir)
      B) Ayrı paragraflarda: "Soru:" ve "Cevap:" satırları ayrı

    Args:
        docx_path: .docx dosyasının yolu.

    Returns:
        Her QA çifti için {question, answer, section} listesi.
    """
    from docx import Document

    docx_path = Path(docx_path)
    if not docx_path.exists():
        raise FileNotFoundError(f"DOCX dosyası bulunamadı: {docx_path}")

    doc = Document(str(docx_path))

    qa_pairs: list[dict] = []
    current_section: str = ""
    current_question: str = ""
    current_answer_lines: list[str] = []

    # \xa0 (non-breaking space) → normal boşluk
    def _clean(s: str) -> str:
        return s.replace("\xa0", " ").strip()

    def _flush_qa() -> None:
        """Biriken soru-cevap çiftini listeye ekle."""
        nonlocal current_question, current_answer_lines
        if current_question and current_answer_lines:
            answer = _normalize("\n".join(current_answer_lines))
            qa_pairs.append({
                "question": _clean(current_question),
                "answer": _clean(answer),
                "section": current_section,
            })
        current_question = ""
        current_answer_lines = []

    for para in doc.paragraphs:
        raw = para.text.strip()
        if not raw:
            continue

        # \xa0 temizle
        text = raw.replace("\xa0", " ")

        # Bölüm başlığı algılama (ör: "1. Yeri Uygulama Eğitimi...")
        section_match = re.match(r"^(\d+)\.\s+(.+)", text)
        if section_match and "Soru:" not in text and "Cevap:" not in text:
            _flush_qa()
            current_section = text
            continue

        # Tek paragrafta Soru + Cevap ("Soru:…\nCevap:…")
        if "Soru:" in text and "Cevap:" in text:
            _flush_qa()
            # \n ile ayrılmış olabilir
            soru_part, *cevap_parts = text.split("Cevap:")
            q = soru_part.replace("Soru:", "").strip()
            a = "Cevap:".join(cevap_parts).strip()  # birden fazla "Cevap:" olabilir
            if q and a:
                qa_pairs.append({
                    "question": q,
                    "answer": a,
                    "section": current_section,
                })
            continue

        # Sadece Soru satırı
        if text.startswith("Soru:"):
            _flush_qa()
            current_question = text[len("Soru:"):].strip()
            continue

        # Sadece Cevap satırı
        if text.startswith("Cevap:"):
            answer_text = text[len("Cevap:"):].strip()
            if answer_text:
                current_answer_lines.append(answer_text)
            continue

        # Devam satırı (mevcut cevabın devamı)
        if current_question:
            current_answer_lines.append(text)

    # Son kalan QA'yı ekle
    _flush_qa()

    logger.info("DOCX'ten %d soru-cevap çifti çıkarıldı: %s", len(qa_pairs), docx_path.name)
    return qa_pairs


# ══════════════════════════════════════════════════════════════════════
#  PARÇALAMA (CHUNKING)
# ══════════════════════════════════════════════════════════════════════

def chunk_text(
    slides: list[SlideText],
    chunk_size: int = 550,
    overlap: int = 80,
    source: str = "pptx",
) -> list[Chunk]:
    """
    Slayt metinlerini ~chunk_size karakter uzunluğunda, overlap kadar
    örtüşen parçalara böler.

    Kısa slaytlar (<= chunk_size) tek parça olarak kalır.
    Uzun slaytlar birden fazla parçaya bölünür.
    """
    chunks: list[Chunk] = []
    idx = 0

    for slide in slides:
        text = slide["text"]
        slide_no = slide["slide_number"]

        if len(text) <= chunk_size:
            chunks.append(Chunk(
                text=text, slide_number=slide_no, chunk_index=idx, source=source,
            ))
            idx += 1
        else:
            start = 0
            while start < len(text):
                end = start + chunk_size
                fragment = text[start:end]

                if end < len(text):
                    last_space = fragment.rfind(" ")
                    if last_space > chunk_size * 0.3:
                        fragment = fragment[:last_space]
                        end = start + last_space

                chunks.append(Chunk(
                    text=fragment.strip(), slide_number=slide_no,
                    chunk_index=idx, source=source,
                ))
                idx += 1
                start = end - overlap if end - overlap > start else end

    logger.info("Toplam %d parça oluşturuldu (%s).", len(chunks), source)
    return chunks


def chunk_docx_qa(
    qa_pairs: list[dict],
    start_index: int = 0,
) -> list[Chunk]:
    """
    DOCX Soru-Cevap çiftlerini chunk olarak döndürür.

    Her QA çifti kendi başına bir chunk olur (SSS formatı zaten kısa).
    """
    chunks: list[Chunk] = []

    for i, qa in enumerate(qa_pairs):
        # Soru + Cevap'ı birleştir
        text = f"Soru: {qa['question']}\nCevap: {qa['answer']}"
        chunks.append(Chunk(
            text=text,
            slide_number=None,
            chunk_index=start_index + i,
            source="docx",
        ))

    logger.info("DOCX'ten %d QA chunk oluşturuldu.", len(chunks))
    return chunks


# ── Tek çağrı fonksiyonları ───────────────────────────────────────────
def load_and_chunk_pptx(
    pptx_path: str | Path,
    chunk_size: int = 550,
    overlap: int = 80,
) -> list[Chunk]:
    """PPTX'ten metin çıkar ve parçalara ayır."""
    slides = extract_slides(pptx_path)
    return chunk_text(slides, chunk_size=chunk_size, overlap=overlap, source="pptx")


def load_and_chunk_docx(
    docx_path: str | Path,
    start_index: int = 0,
) -> list[Chunk]:
    """DOCX SSS dosyasından QA chunk'ları oluştur."""
    qa_pairs = extract_docx_qa(docx_path)
    return chunk_docx_qa(qa_pairs, start_index=start_index)


# Geriye uyumluluk
load_and_chunk = load_and_chunk_pptx
